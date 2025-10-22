import pandas as pd
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

#####################################################################################

def get_tables_as_dataframes(presentation, slide_number):
    """
    Retorna as tabelas de um slide específico como DataFrames.
    """
    try:
        slide = presentation.slides[slide_number]
        tables = []
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                # Converter a tabela para DataFrame
                df = pd.DataFrame(table_data)
                tables.append((shape, df))
        return tables
    except IndexError:
        return None

#####################################################################################

def update_table(shape, edited_df, tamanho_fonte=11):
    """
    Substitui os dados da tabela no shape com os dados do DataFrame editado.
    """
    table = shape.table
    # Ajustar número de linhas e colunas se necessário
    while len(table.rows) < len(edited_df):
        table.add_row()
    while len(table.columns) < len(edited_df.columns):
        table.add_column()
    
    # Atualizar o conteúdo das células
    for i, row in enumerate(edited_df.itertuples(index=False)):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(tamanho_fonte)
                paragraph.alignment = PP_ALIGN.CENTER

#####################################################################################

def substituir_valores(valor, replace_dict, with_pct):
    if isinstance(valor, str):  # Apenas processa valores do tipo string
        for key, new_value in replace_dict.items():
            if isinstance(new_value, (int, float)):
                formatted_value = f"{new_value:.2f}"
            else:
                formatted_value = str(new_value)

            offset = 4
            if with_pct:
                valor = valor.replace(str(key), f"{' ' * offset}{formatted_value}%")
            else:
                valor = valor.replace(str(key), f"{' ' * offset}{formatted_value}")
    
    return valor

#####################################################################################

def substituir_valores_df(df, replace_dict, with_pct):
    return df.applymap(lambda cell: substituir_valores(cell, replace_dict, with_pct))

#####################################################################################
