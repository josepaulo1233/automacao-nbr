from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.util import Inches
from PIL import Image

#####################################################################################

def carregar_pptx(uploaded_file):
    return Presentation(uploaded_file)

#####################################################################################

def extrair_textos(ppt):
    textos = []
    for i, slide in enumerate(ppt.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() != "":
                textos.append((i, shape, shape.text))
    return textos

#####################################################################################

def alterar_texto(ppt, slide_index, old_text, new_text, color_fonte, bold, font_size, hyper_link=False):

    slide = ppt.slides[slide_index]

    if color_fonte == 'white':
        rgb = RGBColor(255, 255, 255)
    else:
        rgb = RGBColor(0, 0, 0)

    for shape in slide.shapes:
        
        if hasattr(shape, 'text') and shape.text == old_text:

            if hyper_link == False:

                # Modificando o texto
                shape.text = new_text

                # Ajustando o tamanho da fonte se necessário
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = Pt(font_size)
                    paragraph.font.bold = bold
                    paragraph.font.color.rgb = rgb

            else:
                shape.text = ""  # Limpar o texto antigo
                text_frame = shape.text_frame
                
                # Adicionar um novo parágrafo com o texto de link
                p = text_frame.add_paragraph()
                p.text = ""  # Inicializar vazio

                # Configurar o texto do link
                run = p.add_run()
                run.text = 'Link entorno'
                run.hyperlink.address = new_text

                # Estilo do texto
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = rgb
                    
    return ppt

#####################################################################################

def adicionar_imagem_com_proporcao(prs, indice_slide, caminho_imagem, largura_util_cm=22.34, altura_max_cm=10.84, margem_topo_cm=6.0, margem_esquerda_cm=0.2):

    # Abre a imagem e obtém as dimensões
    img = Image.open(caminho_imagem)
    largura_px, altura_px = img.size
    proporcao = largura_px / altura_px

    # Converte limites de cm para polegadas
    largura_util_in = largura_util_cm / 2.54
    altura_max_in = altura_max_cm / 2.54
    margem_topo_in = margem_topo_cm / 2.54
    margem_esquerda_in = margem_esquerda_cm / 2.54

    # Calcula dimensões finais mantendo a proporção
    largura_final_in = largura_util_in
    altura_final_in = largura_final_in / proporcao

    if altura_final_in > altura_max_in:
        altura_final_in = altura_max_in
        largura_final_in = altura_final_in * proporcao

    # Calcula posição para centralizar horizontalmente com margem esquerda ajustada
    slide_largura_in = prs.slide_width.inches
    left_in = margem_esquerda_in + (slide_largura_in - largura_final_in - margem_esquerda_in) / 2

    # Aponta para o slide e insere a imagem
    slide = prs.slides[indice_slide]
    slide.shapes.add_picture(
        caminho_imagem,
        Inches(left_in),
        Inches(margem_topo_in),
        width=Inches(largura_final_in),
        height=Inches(altura_final_in)
    )

    return prs

#####################################################################################